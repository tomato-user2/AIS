import pptx
from pptx.util import Inches

def create_presentation():
    # Create a PowerPoint presentation
    presentation = pptx.Presentation()

    # Title Slide
    slide = presentation.slides.add_slide(presentation.slide_layouts[0])
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Machine Learning Project: Navigational Status Prediction"
    subtitle.text = "Team: [Your Team's Name]\nDate: [Insert Date Here]"

    # Introduction Slide 1
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    title.text = "Introduction"
    content.text = (
        "This project aims to predict the navigational status of vessels using AIS data.\n\n"
        "The dataset includes positional and motion data, focusing on binary classification tasks:\n"
        "- 'Under way using engine'\n"
        "- 'Engaged in fishing'.\n\n"
        "We preprocess the data to ensure consistency and generalizability."
    )

    # Introduction Slide 2
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    title.text = "Dataset Overview"
    content.text = (
        "The dataset consists of AIS reports aggregated over time.\n\n"
        "Key steps include:\n"
        "- Grouping data by vessel\n"
        "- Filtering for open sea and anchorage areas\n"
        "- Creating uniform 15-minute and 2-hour sets\n"
        "- Interpolating missing values for columns with <50% missing data."
    )

    # Literature Review Slide
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    title.text = "Literature Review"
    content.text = (
        "We reviewed related work in vessel behavior prediction and motion analysis:\n\n"
        "1. Approaches using positional data for classification.\n"
        "2. Time-series modeling for maritime applications.\n"
        "3. Predictive modeling of vessel trajectories.\n\n"
        "These informed our preprocessing and modeling choices."
    )

    # Dataset Characteristics Slide 1
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    title.text = "Dataset Characteristics"
    content.text = (
        "The dataset includes:\n\n"
        "- Timestamp\n"
        "- Navigational status (target)\n"
        "- Speed over ground (SOG)\n"
        "- Course over ground (COG)\n"
        "- Heading\n\n"
        "These features enable movement-based predictions without positional bias."
    )

    # Dataset Characteristics Slide 2
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    title.text = "Preprocessing"
    content.text = (
        "Preprocessing steps include:\n\n"
        "- Normalizing each feature\n"
        "- Subtracting initial course to standardize direction\n"
        "- Creating consistent time intervals\n"
        "- Removing unnecessary columns to reduce noise."
    )

    # Baseline Model Slide
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    title.text = "Baseline Model"
    content.text = (
        "We tested logistic regression and random forest models:\n\n"
        "- Logistic regression: Simple but limited accuracy (0.68-0.78).\n"
        "- Random forest: Improved performance but struggled with class imbalance.\n\n"
        "Feature importance analysis helped refine preprocessing."
    )

    # Model Definition and Evaluation Slide 1
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    title.text = "Model Definition"
    content.text = (
        "We implemented an LSTM model for time-series classification:\n\n"
        "- Inputs: 15-minute sequences of SOG, COG, and Heading.\n"
        "- Outputs: Binary classification of navigational status.\n\n"
        "Normalization and cyclical encoding were key preprocessing steps."
    )

    # Model Definition and Evaluation Slide 2
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    title.text = "Evaluation"
    content.text = (
        "Evaluation metrics:\n\n"
        "- Accuracy: 0.85 on the best model.\n"
        "- Confusion matrix analysis revealed improved class balance.\n\n"
        "Hyperparameter tuning and sliding window approaches enhanced performance."
    )

    # Results Slide 1
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    title.text = "Results"
    content.text = (
        "Key results:\n\n"
        "- Baseline models: Accuracy of 0.68-0.78.\n"
        "- LSTM models: Accuracy of 0.85 with improved class recognition.\n\n"
        "Visualization of trajectories highlighted model insights."
    )

    # Results Slide 2
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    title.text = "Visualization"
    content.text = (
        "Color-coded trajectory maps provided additional insights:\n\n"
        "- Speed variations were clearly visible.\n"
        "- Differences between navigational statuses were more apparent."
    )

    # Challenges and Errors Slide
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    title.text = "Challenges and Errors"
    content.text = (
        "Major challenges:\n\n"
        "- Data preprocessing and standardization were time-consuming.\n"
        "- Class imbalance affected early models.\n"
        "- Hyperparameter tuning provided limited improvements.\n\n"
        "Refining preprocessing and evaluation strategies resolved many issues."
    )

    # Discussion Slide
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    title.text = "Discussion"
    content.text = (
        "Reflections and future work:\n\n"
        "- Improved preprocessing and visualization enhanced model insights.\n"
        "- Future work includes dynamic timeslot modeling and additional features like wind data.\n"
        "- Sliding window approaches show promise for real-time applications."
    )

    # Conclusion and Future Work Slide
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    title.text = "Conclusion and Future Work"
    content.text = (
        "Conclusions:\n\n"
        "- The LSTM model achieved significant improvements over baseline models.\n"
        "- Visualization techniques provided deeper insights into vessel behavior.\n\n"
        "Future directions:\n"
        "- Expanding dataset diversity.\n"
        "- Incorporating additional features.\n"
        "- Developing dynamic prediction models."
    )

    # Q&A Slide
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    title.text = "Q&A"
    content.text = "Thank you for your attention! We welcome your questions."

    # Save the presentation
    presentation.save("4_Presentation/ML_Project_Presentation.pptx")

if __name__ == "__main__":
    create_presentation()
